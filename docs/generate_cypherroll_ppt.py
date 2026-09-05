#!/usr/bin/env python3
"""Generate CypherRoll deep-dive PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml
from lxml import etree
import copy

# Brand colors
BG = RGBColor(0x0F, 0x17, 0x2A)
CARD = RGBColor(0x1E, 0x29, 0x3B)
GOLD = RGBColor(0xF5, 0x9E, 0x0B)
AMBER = RGBColor(0xFB, 0xBF, 0x24)
VIOLET = RGBColor(0x8B, 0x5C, 0xF6)
GREEN = RGBColor(0x10, 0xB9, 0x81)
RED = RGBColor(0xEF, 0x44, 0x44)
WHITE = RGBColor(0xF8, 0xFA, 0xFC)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
SLATE = RGBColor(0x33, 0x41, 0x55)
DARK = RGBColor(0x02, 0x06, 0x17)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def set_slide_bg(slide, color=BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    # send to back
    spTree = slide.shapes._spTree
    sp = bg._element
    spTree.remove(sp)
    spTree.insert(2, sp)
    return bg


def add_rect(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1)
    else:
        s.line.fill.background()
    # tighten corner radius via adj
    try:
        s.adjustments[0] = 0.08
    except Exception:
        pass
    return s


def add_sharp_rect(slide, l, t, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    return s


def add_text(slide, l, t, w, h, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return box


def add_bullets(slide, l, t, w, h, items, size=14, color=WHITE, bold_first=False, spacing=6):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(spacing)
        p.level = 0
        if bold_first and i == 0:
            p.font.bold = True
    return box


def add_accent_bar(slide, color=GOLD):
    add_sharp_rect(slide, Inches(0), Inches(0), Inches(0.12), prs.slide_height, color)


def add_footer(slide, num, total=28):
    add_text(slide, Inches(0.4), Inches(7.1), Inches(6), Inches(0.3),
             "CypherRoll  ·  Technical Deep Dive", size=10, color=MUTED)
    add_text(slide, Inches(11.5), Inches(7.1), Inches(1.5), Inches(0.3),
             f"{num} / {total}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def section_header(slide, title, subtitle=None):
    add_accent_bar(slide)
    add_text(slide, Inches(0.5), Inches(0.25), Inches(12), Inches(0.5),
             title, size=28, bold=True, color=GOLD)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.75), Inches(12), Inches(0.35),
                 subtitle, size=14, color=MUTED)


TOTAL = 28

# ═══════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
add_sharp_rect(s, 0, Inches(2.4), prs.slide_width, Inches(0.08), GOLD)
add_sharp_rect(s, 0, Inches(4.9), prs.slide_width, Inches(0.04), VIOLET)

add_text(s, Inches(0.8), Inches(1.4), Inches(11.5), Inches(0.4),
         "AUTONOMOUS  ·  PROVABLY FAIR  ·  PRIVACY-FIRST", size=14, color=VIOLET, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(2.6), Inches(11.5), Inches(0.9),
         "CYPHERROLL", size=60, bold=True, color=GOLD, align=PP_ALIGN.CENTER, font="Arial")
add_text(s, Inches(0.8), Inches(3.5), Inches(11.5), Inches(0.5),
         "The Next-Generation Provably Fair, High-Performance\nWeb3 & Privacy-First Casino",
         size=18, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.4),
         "Technical Deep-Dive Presentation  ·  Architecture · Cryptography · Games · Vaults · Ops",
         size=13, color=MUTED, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.3),
         "Sub-50ms latency  ·  HMAC-SHA256 commit–reveal  ·  Multi-chain escrow  ·  Tor v3",
         size=12, color=AMBER, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
section_header(s, "Agenda", "28-slide deep dive — nothing skipped")
add_footer(s, 2)

cols = [
    ("01–05  Foundation", [
        "• Vision & problem space",
        "• Value proposition",
        "• Audience & personas",
        "• Architecture overview",
        "• Full tech stack",
    ]),
    ("06–12  Product Core", [
        "• CypherDice mechanics",
        "• CypherCrash engine",
        "• Provably fair crypto",
        "• Exact formulas",
        "• Wallet auth (SIWE/SIWS)",
        "• Smart contracts",
        "• Cashier flows",
    ]),
    ("13–20  Systems", [
        "• Treasury & Kelly caps",
        "• VIP rakeback economy",
        "• Design system & 3D",
        "• Social (ticker/trollbox)",
        "• Privacy (Tor/PoW/AML)",
        "• Admin command center",
        "• API & data model",
        "• Security invariants",
    ]),
    ("21–28  Delivery", [
        "• Deployment (Docker/Tor)",
        "• Demo vs Real mode",
        "• Status vs roadmap",
        "• End-to-end user journeys",
        "• Verification checklist",
        "• Key takeaways",
        "• Appendix formulas",
        "• Closing",
    ]),
]
for i, (title, items) in enumerate(cols):
    x = Inches(0.4 + i * 3.2)
    add_rect(s, x, Inches(1.3), Inches(3.0), Inches(5.4), CARD, SLATE)
    add_text(s, x + Inches(0.15), Inches(1.45), Inches(2.7), Inches(0.4),
             title, size=13, bold=True, color=GOLD)
    add_bullets(s, x + Inches(0.15), Inches(2.0), Inches(2.7), Inches(4.5),
                items, size=12, color=WHITE, spacing=8)

# ═══════════════════════════════════════════════════════════
# SLIDE 3 — Problem
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
section_header(s, "The Problem Space", "Why opaque casinos fail crypto-native users")
add_footer(s, 3)

problems = [
    ("Opaque RNG", RED, "Players cannot verify outcomes. House can alter results after bets. Trust is social, not mathematical."),
    ("On-chain Latency", VIOLET, "Waiting for block confirmations kills high-frequency betting. Micro-bets become unusable."),
    ("KYC Friction", GOLD, "Email/password + documents contradict Web3 identity. Privacy advocates refuse to play."),
    ("Solvency Opacity", AMBER, "No proof the house can pay. Bankroll risk is hidden until withdrawal fails."),
]
for i, (title, color, desc) in enumerate(problems):
    y = Inches(1.35 + i * 1.3)
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(1.15), CARD, SLATE)
    add_sharp_rect(s, Inches(0.5), y, Inches(0.12), Inches(1.15), color)
    add_text(s, Inches(0.85), y + Inches(0.15), Inches(11.5), Inches(0.35),
             title, size=18, bold=True, color=color)
    add_text(s, Inches(0.85), y + Inches(0.55), Inches(11.5), Inches(0.45),
             desc, size=14, color=WHITE)

# ═══════════════════════════════════════════════════════════
# SLIDE 4 — Vision & Value
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
section_header(s, "Vision & Value Proposition",
               "“Autonomous Provably Fair Gaming” — Multi-Chain Cryptographic Sessions • SIWE & SIWS Verified")
add_footer(s, 4)

props = [
    ("Provably Fair", "HMAC-SHA256 commit–reveal. Player audits every roll. Server seed hashed before play."),
    ("Sub-50ms Speed", "Off-chain game ticks. On-chain only for deposits/withdrawals. Instant resolution."),
    ("Zero-KYC Auth", "Wallet connect + signed message. No email, no password, no documents."),
    ("Privacy-First", "Tor Onion v3, WebRTC leak mitigation, zero IP logging, nginx access_log off."),
    ("Non-Custodial", "EVM + Solana escrow vaults. Operator-signed EIP-712 withdrawals. Player controls keys."),
    ("Solvency Aware", "Community Bankroll Vault + Kelly 1% caps + live HEALTHY / WARNING / CIRCUIT_BREAKER."),
]
for i, (title, desc) in enumerate(props):
    col, row = i % 3, i // 3
    x, y = Inches(0.4 + col * 4.25), Inches(1.35 + row * 2.6)
    add_rect(s, x, y, Inches(4.05), Inches(2.35), CARD, SLATE)
    add_text(s, x + Inches(0.2), y + Inches(0.25), Inches(3.6), Inches(0.4),
             title, size=16, bold=True, color=GOLD)
    add_text(s, x + Inches(0.2), y + Inches(0.8), Inches(3.6), Inches(1.3),
             desc, size=13, color=WHITE)

# ═══════════════════════════════════════════════════════════
# SLIDE 5 — Audience
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
section_header(s, "Target Audience & User Stories", "Five personas from the PRD")
add_footer(s, 5)

stories = [
    ("Web3 Bettor", "Connect Phantom / MetaMask with one click (SIWE/SIWS). Play immediately — no registration forms."),
    ("High-Speed Gambler", "Instantaneous resolution (<50ms) in Crash & Dice. Micro-bets without waiting for blocks."),
    ("Fairness Skeptic", "Inspect SHA-256 pre-committed Server Seed, set Client Seed, verify HMAC formula independently."),
    ("Privacy Advocate", "Access via Tor .onion with zero WebRTC IP leakage. Identity & location stay confidential."),
    ("LP / Investor", "Stake USDC into Community Bankroll Vault. Earn proportional share of house-edge revenue."),
]
for i, (title, desc) in enumerate(stories):
    y = Inches(1.25 + i * 1.1)
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.95), CARD, SLATE)
    add_text(s, Inches(0.7), y + Inches(0.12), Inches(3.2), Inches(0.35),
             f"As a {title}", size=14, bold=True, color=AMBER)
    add_text(s, Inches(4.0), y + Inches(0.2), Inches(8.5), Inches(0.6),
             desc, size=13, color=WHITE)

# ═══════════════════════════════════════════════════════════
# SLIDE 6 — Architecture
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
section_header(s, "Architecture — 3-Tier Hybrid Web3",
               "Off-Chain Execution + On-Chain Settlement")
add_footer(s, 6)

tiers = [
    (GOLD, "PRESENTATION", [
        "Next.js 14 App Router",
        "Tailwind + Framer Motion",
        "React Three Fiber 3D",
        "RainbowKit + Solana adapters",
        "Lobby tabs: DICE | CRASH | VAULT",
    ]),
    (VIOLET, "APPLICATION", [
        "Game engines (Dice, Crash)",
        "HMAC provably-fair core",
        "SIWE / SIWS auth + sessions",
        "Kelly caps · EIP-712 signer",
        "PoW · AML · Admin controls",
    ]),
    (GREEN, "PERSISTENCE", [
        "Supabase Postgres + RLS",
        "Atomic RPCs (FOR UPDATE)",
        "Optional Redis / Upstash",
        "Realtime (bets, trollbox)",
        "In-memory fallback mock DB",
    ]),
    (AMBER, "ON-CHAIN", [
        "CypherRollVault (EVM)",
        "CypherRollVRFConsumer",
        "Solana cypherroll_escrow",
        "Base 8453 · Arbitrum 42161",
        "USDC deposits / withdrawals",
    ]),
]
for i, (color, title, items) in enumerate(tiers):
    x = Inches(0.35 + i * 3.25)
    add_rect(s, x, Inches(1.3), Inches(3.1), Inches(5.4), CARD, SLATE)
    add_sharp_rect(s, x, Inches(1.3), Inches(3.1), Inches(0.5), color)
    add_text(s, x, Inches(1.38), Inches(3.1), Inches(0.4),
             title, size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    add_bullets(s, x + Inches(0.2), Inches(2.0), Inches(2.7), Inches(4.4),
                [f"• {x}" for x in items], size=13, color=WHITE, spacing=10)

# ═══════════════════════════════════════════════════════════
# SLIDE 7 — Tech Stack
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
section_header(s, "Technology Stack", "Exact versions from package.json")
add_footer(s, 7)

stack = [
    ("Frontend", [
        "Next.js ^14.2.3 (App Router)",
        "React ^18.3.1 · TypeScript ^5.4.5",
        "Tailwind CSS ^3.4.3",
        "Framer Motion ^11.2.6",
        "three ^0.164.1 (R3F canvases)",
        "lucide-react · clsx · tw-merge",
    ]),
    ("Web3", [
        "wagmi ^2.8.6 · viem ^2.10.8",
        "RainbowKit ^2.1.2",
        "@solana/web3.js ^1.91.8",
        "wallet-adapter 0.9–0.19",
        "Solidity ^0.8.24 / ^0.8.20",
        "Anchor (Solana escrow)",
    ]),
    ("Data & Runtime", [
        "@supabase/supabase-js ^2.115",
        "@tanstack/react-query ^5.37",
        "Redis / Upstash (optional)",
        "Node 20 Alpine Docker",
        "nginx + Tor Onion v3",
        "Blender headless (3D gen)",
    ]),
    ("Contracts / Assets", [
        "CypherRollVault.sol",
        "CypherRollVRFConsumer.sol",
        "cypherroll_escrow (Rust)",
        "dice.glb · rocket.glb · chip.glb",
        "npm run test:pf (PF suite)",
        "npm run generate:3d",
    ]),
]
for i, (title, items) in enumerate(stack):
    x = Inches(0.35 + i * 3.25)
    add_rect(s, x, Inches(1.3), Inches(3.1), Inches(5.4), CARD, SLATE)
    add_text(s, x + Inches(0.15), Inches(1.45), Inches(2.8), Inches(0.4),
             title, size=16, bold=True, color=GOLD)
    add_bullets(s, x + Inches(0.15), Inches(2.05), Inches(2.8), Inches(4.4),
                [f"• {x}" for x in items], size=12, color=WHITE, spacing=9)

# ═══════════════════════════════════════════════════════════
# SLIDE 8 — CypherDice
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
section_header(s, "CypherDice — Deep Dive", "Roll-under · 1% house edge · 99% RTP · 3D canvas")
add_footer(s, 8)

add_rect(s, Inches(0.4), Inches(1.25), Inches(6.2), Inches(5.5), CARD, SLATE)
add_text(s, Inches(0.6), Inches(1.4), Inches(5.8), Inches(0.35),
         "Mechanics", size=16, bold=True, color=GOLD)
add_bullets(s, Inches(0.6), Inches(1.9), Inches(5.8), Inches(4.6), [
    "• Target range: roll-under 2–98 (UI slider)",
    "• Outcome space: 0.00 – 99.99 (two decimals)",
    "• Win condition: roll < target",
    "• Multiplier: (1 − 0.01) / (target / 100)",
    "• House edge fixed at 1.00% → RTP 99%",
    "• Demo mode: virtual $1,000 (localStorage)",
    "• Real mode: session + atomic DB bet RPC",
    "• History: each bet opens in PF verifier",
    "• Rakeback: 15% of theoretical 1% edge",
    "  → rakeback = wager × 0.01 × 0.15",
    "• API: POST /api/games/dice/roll",
    "• 3D: DiceCanvas (dice.glb ~267KB) + 2D fallback",
], size=13, color=WHITE, spacing=6)

add_rect(s, Inches(6.85), Inches(1.25), Inches(6.0), Inches(5.5), CARD, SLATE)
add_text(s, Inches(7.05), Inches(1.4), Inches(5.6), Inches(0.35),
         "Outcome Pipeline", size=16, bold=True, color=VIOLET)
add_bullets(s, Inches(7.05), Inches(1.9), Inches(5.6), Inches(4.6), [
    "1. CSPRNG serverSeed = randomBytes(32)",
    "2. Publish SHA-256(serverSeed) BEFORE play",
    "3. Player sets clientSeed; nonce++",
    "4. HMAC = HMAC-SHA256(serverSeed,",
    "      `${clientSeed}:${nonce}`)",
    "5. subHash = HMAC[0:8]  (32 bits)",
    "6. roll = (parseInt(subHash,16) % 10000) / 100",
    "7. won = roll < target",
    "8. payout = wager × getDiceMultiplier(target)",
    "",
    "Example: target=50 → mult ≈ 1.9800×",
    "Example: target=25 → mult ≈ 3.9600×",
    "Example: target=2  → mult ≈ 49.5000×",
], size=13, color=WHITE, spacing=5)

# ═══════════════════════════════════════════════════════════
# SLIDE 9 — CypherCrash
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
section_header(s, "CypherCrash — Global Multiplayer Engine",
               "50ms tick · Bustabit-class formula · 2% house edge · 98% RTP")
add_footer(s, 9)

add_rect(s, Inches(0.4), Inches(1.25), Inches(6.2), Inches(5.5), CARD, SLATE)
add_text(s, Inches(0.6), Inches(1.4), Inches(5.8), Inches(0.35),
         "Round Lifecycle", size=16, bold=True, color=GOLD)
add_bullets(s, Inches(0.6), Inches(1.9), Inches(5.8), Inches(4.6), [
    "• States: STARTING (5s) → FLYING → CRASHED (~3.5s)",
    "• Flight curve: mult(t) = exp(0.065 × t_sec)",
    "• Tick interval: 50ms (sub-50ms UX target)",
    "• Bets only accepted in STARTING",
    "• One bet per wallet per round",
    "• Auto-cashout range: 1.01× – 10,000×",
    "• Manual cashout during FLYING",
    "• 250ms latency grace after crash for late packets",
    "• Singleton CrashEngine + leader election",
    "• Distributed lock on bet placement",
    "• SSE stream: /api/games/crash/stream",
    "• Also: /state, /ping, /bet, /cashout",
    "• History: last ~10 crash points",
    "• 3D: CrashRocketCanvas (rocket.glb ~56KB)",
], size=12, color=WHITE, spacing=5)

add_rect(s, Inches(6.85), Inches(1.25), Inches(6.0), Inches(5.5), CARD, SLATE)
add_text(s, Inches(7.05), Inches(1.4), Inches(5.6), Inches(0.35),
         "Crash Point Formula", size=16, bold=True, color=RED)
add_bullets(s, Inches(7.05), Inches(1.9), Inches(5.6), Inches(4.6), [
    "HMAC = HMAC-SHA256(serverSeed,",
    "  `${clientSeed}:${nonce}`)",
    "h = parseInt(HMAC[0:13], 16)   // 52 bits",
    "e = 2^52",
    "raw = (0.98 × e) / (e − h)",
    "crash = raw < 1 ? 1.00",
    "      : floor(raw × 100) / 100",
    "",
    "• Coefficient 0.98 ⇒ ~2% mass at 1.00×",
    "  (instant crash = house edge)",
    "• Industry standard (Bustabit / Rollbit)",
    "• Rakeback on settle:",
    "  wager × 0.02 × 0.15",
    "• PF suite: 10k sims → ~2% instant crashes",
    "• Seeds revealed after round in crash_rounds",
], size=12, color=WHITE, spacing=5)

# ═══════════════════════════════════════════════════════════
# SLIDE 10 — Provably Fair
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
section_header(s, "Provably Fair Cryptography",
               "Commit–reveal · Independent auditor · Optional Chainlink VRF path")
add_footer(s, 10)

steps = [
    ("01  COMMIT", GOLD, "Generate serverSeed via crypto.randomBytes(32). Publish only SHA-256(serverSeed) before any bet."),
    ("02  PLAY", VIOLET, "Player chooses clientSeed. Nonce increments per bet. Outcome = f(HMAC(serverSeed, clientSeed:nonce))."),
    ("03  REVEAL", GREEN, "On seed rotation (or crash round end), reveal serverSeed. Anyone recomputes hash + HMAC."),
    ("04  VERIFY", AMBER, "In-app ProvablyFairModal + POST /api/games/provably-fair/verify. getDetailedGameProof() returns step-by-step."),
]
for i, (title, color, desc) in enumerate(steps):
    y = Inches(1.25 + i * 1.25)
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(1.1), CARD, SLATE)
    add_sharp_rect(s, Inches(0.5), y, Inches(0.12), Inches(1.1), color)
    add_text(s, Inches(0.85), y + Inches(0.15), Inches(11.5), Inches(0.3),
             title, size=15, bold=True, color=color)
    add_text(s, Inches(0.85), y + Inches(0.5), Inches(11.5), Inches(0.45),
             desc, size=13, color=WHITE)

# ═══════════════════════════════════════════════════════════
# SLIDE 11 — Formulas appendix visual
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
section_header(s, "Exact Formulas — Side by Side", "Source of truth: lib/provably-fair.ts")
add_footer(s, 11)

add_rect(s, Inches(0.4), Inches(1.25), Inches(6.2), Inches(5.5), CARD, SLATE)
add_text(s, Inches(0.6), Inches(1.4), Inches(5.8), Inches(0.35),
         "DICE", size=18, bold=True, color=GOLD)
add_bullets(s, Inches(0.6), Inches(1.95), Inches(5.8), Inches(4.5), [
    "computeHMAC(serverSeed, clientSeed, nonce)",
    "  → HMAC-SHA256 key=serverSeed",
    "  → msg = `${clientSeed}:${nonce}`",
    "",
    "calculateDiceRoll(...)",
    "  subHash = hex[0:8]",
    "  roll = (int(subHash,16) % 10000) / 100",
    "  → [0.00, 99.99]",
    "",
    "getDiceMultiplier(target)",
    "  target ∈ [1.00, 98.00]",
    "  mult = (1 − 0.01) / (target/100)",
    "  → 4 decimal places",
    "",
    "Invariant: NEVER Math.random() for outcomes",
], size=12, color=WHITE, spacing=4)

add_rect(s, Inches(6.85), Inches(1.25), Inches(6.0), Inches(5.5), CARD, SLATE)
add_text(s, Inches(7.05), Inches(1.4), Inches(5.6), Inches(0.35),
         "CRASH + VRF MIX", size=18, bold=True, color=VIOLET)
add_bullets(s, Inches(7.05), Inches(1.95), Inches(5.6), Inches(4.5), [
    "calculateCrashPoint(...)",
    "  h = int(hex[0:13], 16)  // 52-bit",
    "  e = 2^52",
    "  raw = (0.98 * e) / (e - h)",
    "  clamp floor to 2 decimals",
    "",
    "Flight (visual, not RNG)",
    "  multiplier(t) = e^(0.065 * t)",
    "",
    "Optional VRF mix path",
    "  mixed = SHA256(`${word}:${client}:${nonce}`)",
    "  outcome = (int(mixed[0:8],16)/0xffffffff)*100",
    "",
    "Contract: CypherRollVRFConsumer",
    "  VRF v2.5 · 3 conf · 200k gas · 1 word",
    "  TS path currently simulates fulfillment",
], size=12, color=WHITE, spacing=4)

# ═══════════════════════════════════════════════════════════
# SLIDE 12 — Auth
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
section_header(s, "Wallet Auth — SIWE / SIWS", "Zero-KYC cryptographic sessions")
add_footer(s, 12)

flow = [
    ("1. Connect", "Navbar: SOL or EVM. Phantom / Solflare or MetaMask / Rabby / Coinbase via RainbowKit."),
    ("2. Nonce", "GET /api/auth/nonce → challenge: “Welcome to CypherRoll Casino!…”"),
    ("3. Sign", "Wallet signs SIWE-style (EVM) or SIWS-style (Solana) message."),
    ("4. Verify", "POST /api/auth/verify → HMAC-signed cypher_session cookie (TTL 365 days)."),
    ("5. Profile", "Auto-provision: seed pair, Bronze VIP, balance_usdc = 0."),
    ("6. Session", "GET /api/auth/me · POST /api/auth/logout. Spoofed wallets rejected."),
]
for i, (title, desc) in enumerate(flow):
    col, row = i % 3, i // 3
    x, y = Inches(0.4 + col * 4.25), Inches(1.3 + row * 2.7)
    add_rect(s, x, y, Inches(4.05), Inches(2.4), CARD, SLATE)
    add_text(s, x + Inches(0.2), y + Inches(0.25), Inches(3.6), Inches(0.4),
             title, size=16, bold=True, color=GOLD)
    add_text(s, x + Inches(0.2), y + Inches(0.8), Inches(3.6), Inches(1.3),
             desc, size=13, color=WHITE)

# ═══════════════════════════════════════════════════════════
# SLIDE 13 — Smart Contracts
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
section_header(s, "Smart Contracts & Chains", "EVM vaults + Solana escrow + VRF consumer")
add_footer(s, 13)

add_rect(s, Inches(0.4), Inches(1.25), Inches(4.1), Inches(5.5), CARD, SLATE)
add_text(s, Inches(0.55), Inches(1.4), Inches(3.8), Inches(0.4),
         "CypherRollVault.sol", size=14, bold=True, color=GOLD)
add_bullets(s, Inches(0.55), Inches(1.95), Inches(3.8), Inches(4.5), [
    "• depositETH()",
    "• depositERC20(token, amount)",
    "• Withdrawal via EIP-712:",
    "  (player, token, amount,",
    "   nonce, deadline)",
    "• Signed by operatorSigner",
    "• Sequential nonces · pause",
    "• CEI before transfer",
    "• Domain: CypherRollVault v1",
    "• Chains: Base 8453,",
    "  Arbitrum One 42161",
    "• USDC: Base 0x833589…",
    "  Arb 0xaf88d0…",
], size=12, color=WHITE, spacing=4)

add_rect(s, Inches(4.7), Inches(1.25), Inches(4.1), Inches(5.5), CARD, SLATE)
add_text(s, Inches(4.85), Inches(1.4), Inches(3.8), Inches(0.4),
         "VRF + Solana Escrow", size=14, bold=True, color=VIOLET)
add_bullets(s, Inches(4.85), Inches(1.95), Inches(3.8), Inches(4.5), [
    "CypherRollVRFConsumer.sol",
    "• Chainlink VRF v2.5",
    "• requestEntropy(roundId)",
    "• rawFulfillRandomWords",
    "• Base coordinator",
    "  0xd5D517aBE5cF79B7…",
    "",
    "cypherroll_escrow (Anchor)",
    "• initialize_vault",
    "• deposit_sol / withdraw_sol",
    "• PDA: deposited + nonce",
    "• Program id placeholder",
    "• Frontend RPC: Solana devnet",
], size=12, color=WHITE, spacing=4)

add_rect(s, Inches(9.0), Inches(1.25), Inches(3.9), Inches(5.5), CARD, SLATE)
add_text(s, Inches(9.15), Inches(1.4), Inches(3.6), Inches(0.4),
         "Settlement Model", size=14, bold=True, color=GREEN)
add_bullets(s, Inches(9.15), Inches(1.95), Inches(3.6), Inches(4.5), [
    "OFF-CHAIN",
    "• Game outcomes (HMAC)",
    "• Balance ledger (Supabase)",
    "• Crash tick engine",
    "",
    "ON-CHAIN",
    "• Deposits into vault",
    "• Signed withdrawals",
    "• Optional VRF entropy",
    "",
    "TOKEN",
    "• Ledger unit: USDC",
    "• Future: $CYPH",
    "  20% net rev buyback/burn",
], size=12, color=WHITE, spacing=4)

# ═══════════════════════════════════════════════════════════
# SLIDE 14 — Cashier
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
section_header(s, "Cashier — Deposit & Withdraw", "Multi-chain USDC · AML · Idempotent credit")
add_footer(s, 14)

add_rect(s, Inches(0.4), Inches(1.25), Inches(6.2), Inches(5.5), CARD, SLATE)
add_text(s, Inches(0.6), Inches(1.4), Inches(5.8), Inches(0.35),
         "Deposit Flow", size=16, bold=True, color=GREEN)
add_bullets(s, Inches(0.6), Inches(1.95), Inches(5.8), Inches(4.5), [
    "1. Open CashierModal (Base / Arb / Sol)",
    "2. Transfer USDC to vault on-chain",
    "3. POST /api/cashier/deposit (intent)",
    "4. POST /api/cashier/verify-deposit",
    "   • tx-hash idempotency check",
    "   • AML screening (sanctions registry)",
    "   • Quarantine if flagged",
    "5. Credit balance_usdc in profiles",
    "6. Ledger row in transactions",
    "",
    "Dev: ALLOW_SIMULATED_DEPOSITS for testing",
    "Listener: lib/web3/deposit-listener.ts",
], size=13, color=WHITE, spacing=5)

add_rect(s, Inches(6.85), Inches(1.25), Inches(6.0), Inches(5.5), CARD, SLATE)
add_text(s, Inches(7.05), Inches(1.4), Inches(5.6), Inches(0.35),
         "Withdraw Flow", size=16, bold=True, color=AMBER)
add_bullets(s, Inches(7.05), Inches(1.95), Inches(5.6), Inches(4.5), [
    "1. Request withdraw (min $10, max $25,000)",
    "2. POST /api/cashier/withdraw",
    "3. Deduct balance atomically",
    "4. EVM: EIP-712 operator voucher signed",
    "   (OPERATOR_SIGNER_PRIVATE_KEY)",
    "5. Player claims on-chain with signature",
    "6. Solana: escrow withdraw params",
    "   (Ed25519 verify noted for production)",
    "",
    "Safety rails",
    "• Max automated withdrawal: 25,000 USDC",
    "• Kelly max-bet still governs gameplay",
    "• Circuit breaker can freeze flows",
], size=13, color=WHITE, spacing=5)

# ═══════════════════════════════════════════════════════════
# SLIDE 15 — Treasury / Kelly
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
section_header(s, "Treasury, Kelly Caps & Solvency",
               "Community Bankroll Vault · Risk of ruin control")
add_footer(s, 15)

add_rect(s, Inches(0.4), Inches(1.25), Inches(4.1), Inches(5.5), CARD, SLATE)
add_text(s, Inches(0.55), Inches(1.4), Inches(3.8), Inches(0.4),
         "Bankroll Vault UI", size=14, bold=True, color=GOLD)
add_bullets(s, Inches(0.55), Inches(1.95), Inches(3.8), Inches(4.5), [
    "• Stake USDC as LP",
    "• Displayed APY: 19.4%",
    "• Kelly max-bet shown live",
    "• Health statuses:",
    "  HEALTHY",
    "  WARNING",
    "  CIRCUIT_BREAKER_ACTIVE",
    "• API:",
    "  GET /treasury/status",
    "  POST /treasury/stake",
    "• Table: bankroll_stakes",
], size=13, color=WHITE, spacing=6)

add_rect(s, Inches(4.7), Inches(1.25), Inches(4.1), Inches(5.5), CARD, SLATE)
add_text(s, Inches(4.85), Inches(1.4), Inches(3.8), Inches(0.4),
         "Reserves (display defaults)", size=14, bold=True, color=VIOLET)
add_bullets(s, Inches(4.85), Inches(1.95), Inches(3.8), Inches(4.5), [
    "• Base:     $650,000",
    "• Arbitrum: $400,000",
    "• Solana:   $200,000",
    "• Total:    $1,250,000",
    "",
    "Kelly rule",
    "  maxBet = 1% × liquid reserves",
    "  → $12,500 at $1.25M",
    "",
    "Invariant",
    "  Σ user balances ≤ vault reserves",
    "",
    "Treasury tiers (docs)",
    "  Hot ≤5% · Warm 15% · Cold 80%",
], size=13, color=WHITE, spacing=5)

add_rect(s, Inches(9.0), Inches(1.25), Inches(3.9), Inches(5.5), CARD, SLATE)
add_text(s, Inches(9.15), Inches(1.4), Inches(3.6), Inches(0.4),
         "Why Kelly?", size=14, bold=True, color=GREEN)
add_bullets(s, Inches(9.15), Inches(1.95), Inches(3.6), Inches(4.5), [
    "• Caps single-outcome",
    "  exposure to 1% of",
    "  liquid hot reserves",
    "• Prevents one whale",
    "  from bankrupting vault",
    "• Aligns LP risk with",
    "  mathematical edge",
    "• Enforced server-side",
    "  on bet placement",
    "• Pairs with circuit",
    "  breaker for ops halt",
], size=13, color=WHITE, spacing=6)

# ═══════════════════════════════════════════════════════════
# SLIDE 16 — VIP
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
section_header(s, "VIP Rakeback Economy", "Return 10–25% of house edge based on wager volume")
add_footer(s, 16)

tiers = [
    ("Bronze", "$0+", "10%", MUTED),
    ("Silver", "$500+", "12.5%", SLATE),
    ("Gold", "$2,500+", "15%", GOLD),
    ("Platinum", "$10,000+", "20%", AMBER),
    ("Diamond", "$50,000+", "25%", VIOLET),
]
for i, (name, thresh, rate, color) in enumerate(tiers):
    x = Inches(0.4 + i * 2.55)
    add_rect(s, x, Inches(1.4), Inches(2.4), Inches(3.2), CARD, color)
    add_text(s, x, Inches(1.7), Inches(2.4), Inches(0.4),
             name, size=16, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(2.3), Inches(2.4), Inches(0.35),
             thresh, size=13, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(2.9), Inches(2.4), Inches(0.6),
             rate, size=28, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(3.6), Inches(2.4), Inches(0.4),
             "of edge", size=12, color=MUTED, align=PP_ALIGN.CENTER)

add_rect(s, Inches(0.4), Inches(4.9), Inches(12.5), Inches(1.7), CARD, SLATE)
add_bullets(s, Inches(0.6), Inches(5.05), Inches(12), Inches(1.4), [
    "• Runtime engines often apply a fixed 15% of edge (Dice: wager×0.01×0.15 · Crash: wager×0.02×0.15); UI tiers advertise 10–25% of edge by volume.",
    "• Claim unclaimed rakeback into balance via VIPRakebackModal · Accrues on real bets only (not demo).",
    "• Future $CYPH: buyback & burn from 20% of net platform revenue (PRD Priority 2 — not yet implemented).",
], size=13, color=WHITE, spacing=6)

# ═══════════════════════════════════════════════════════════
# SLIDE 17 — Design & 3D
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
section_header(s, "Design System & 3D Visuals", "Dark OLED Cyberpunk / High-Tech Web3 Casino")
add_footer(s, 17)

# color swatches
swatches = [
    ("#0F172A", "Background", BG),
    ("#1E293B", "Card", CARD),
    ("#F59E0B", "Gold", GOLD),
    ("#8B5CF6", "CTA Violet", VIOLET),
    ("#10B981", "Win", GREEN),
    ("#EF4444", "Loss", RED),
]
for i, (hexv, name, color) in enumerate(swatches):
    x = Inches(0.4 + i * 2.15)
    add_sharp_rect(s, x, Inches(1.3), Inches(2.0), Inches(0.7), color)
    add_text(s, x, Inches(2.1), Inches(2.0), Inches(0.3),
             name, size=11, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, Inches(2.35), Inches(2.0), Inches(0.25),
             hexv, size=10, color=MUTED, align=PP_ALIGN.CENTER)

add_rect(s, Inches(0.4), Inches(2.85), Inches(6.2), Inches(3.8), CARD, SLATE)
add_text(s, Inches(0.6), Inches(3.0), Inches(5.8), Inches(0.35),
         "Typography & Effects", size=15, bold=True, color=GOLD)
add_bullets(s, Inches(0.6), Inches(3.5), Inches(5.8), Inches(2.9), [
    "• Headings / multipliers: Orbitron",
    "• Body / controls: Exo 2",
    "• Hashes / seeds / wallets: JetBrains Mono",
    "• Effects: glow-gold, glow-purple, glass-panel",
    "• Icons: lucide-react only (Rules.md)",
    "• WCAG contrast · prefers-reduced-motion",
    "• Responsive: 375px → 1440px",
], size=13, color=WHITE, spacing=6)

add_rect(s, Inches(6.85), Inches(2.85), Inches(6.0), Inches(3.8), CARD, SLATE)
add_text(s, Inches(7.05), Inches(3.0), Inches(5.6), Inches(0.35),
         "3D Assets (Blender → GLB)", size=15, bold=True, color=VIOLET)
add_bullets(s, Inches(7.05), Inches(3.5), Inches(5.6), Inches(2.9), [
    "• dice.glb ~267KB — DiceCanvas",
    "• rocket.glb ~56KB — CrashRocketCanvas",
    "• chip.glb ~172KB — CasinoChipCanvas",
    "• Draco-compressed · under 1.5MB total budget",
    "• npm run generate:3d (headless Blender)",
    "• 2D canvas fallbacks + WebGL error boundary",
    "• Tor-safe: WebGL may fail → graceful degrade",
], size=13, color=WHITE, spacing=6)

# ═══════════════════════════════════════════════════════════
# SLIDE 18 — Social
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
section_header(s, "Live Social Layer", "Realtime bets feed + VIP-colored trollbox")
add_footer(s, 18)

add_rect(s, Inches(0.4), Inches(1.3), Inches(6.2), Inches(5.4), CARD, SLATE)
add_text(s, Inches(0.6), Inches(1.5), Inches(5.8), Inches(0.4),
         "LiveBetsTicker", size=18, bold=True, color=GOLD)
add_bullets(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(4.2), [
    "• Global recent bets stream",
    "• Source: Supabase Realtime and/or",
    "  GET /api/bets/recent",
    "• Surfaces game type, wager, outcome,",
    "  multiplier, wallet snippet",
    "• Builds FOMO + social proof without",
    "  leaving the lobby",
    "• Low-overhead vs full WebSocket lobby",
], size=14, color=WHITE, spacing=8)

add_rect(s, Inches(6.85), Inches(1.3), Inches(6.0), Inches(5.4), CARD, SLATE)
add_text(s, Inches(7.05), Inches(1.5), Inches(5.6), Inches(0.4),
         "Trollbox", size=18, bold=True, color=VIOLET)
add_bullets(s, Inches(7.05), Inches(2.2), Inches(5.6), Inches(4.2), [
    "• GET/POST /api/trollbox",
    "• Table: trollbox_messages",
    "• VIP-colored badges by tier",
    "• Presence-style “online” UX",
    "• Toggleable from lobby chrome",
    "• RLS: public read; writes via service_role",
    "• Completes Rollbit-class social feel",
], size=14, color=WHITE, spacing=8)

# ═══════════════════════════════════════════════════════════
# SLIDE 19 — Privacy / Security
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
section_header(s, "Privacy, PoW & AML", "Tor v3 · WebRTC mitigation · bot resistance · sanctions")
add_footer(s, 19)

blocks = [
    ("Tor & Privacy", GOLD, [
        "Docker: app → nginx:8080 → Tor HS :80",
        "Onion v3 hidden service",
        "nginx access_log off",
        "CSP + Permissions-Policy headers",
        "WebRTC disabled messaging",
        "Tor-safe localStorage polyfill",
        "Zero IP logging design goal",
    ]),
    ("Proof of Work", VIOLET, [
        "Minimal PoW for Tor bot resistance",
        "~3 leading hex zeros on SHA256",
        "Challenge HMAC-signed",
        "~60s TTL",
        "API: /api/security/pow",
        "Client: lib/security/pow-client.ts",
        "Secret: POW_SECRET",
    ]),
    ("AML Screening", RED, [
        "Registry: Tornado Cash, Lazarus, …",
        "POST /api/security/aml-check",
        "Flagged deposits → quarantine",
        "Table: aml_sanctions_quarantine",
        "Runs on verify-deposit path",
        "Admin visibility for ops",
        "Complements solvency controls",
    ]),
]
for i, (title, color, items) in enumerate(blocks):
    x = Inches(0.35 + i * 4.3)
    add_rect(s, x, Inches(1.3), Inches(4.1), Inches(5.4), CARD, SLATE)
    add_text(s, x + Inches(0.2), Inches(1.45), Inches(3.7), Inches(0.4),
             title, size=16, bold=True, color=color)
    add_bullets(s, x + Inches(0.2), Inches(2.05), Inches(3.7), Inches(4.4),
                [f"• {x}" for x in items], size=13, color=WHITE, spacing=8)

# ═══════════════════════════════════════════════════════════
# SLIDE 20 — Admin
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
section_header(s, "Admin Operator Command Center", "/admin — wallet-gated CYPHER_OPERATOR")
add_footer(s, 20)

add_rect(s, Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.4), CARD, SLATE)
add_bullets(s, Inches(0.7), Inches(1.55), Inches(12), Inches(5.0), [
    "• Auth: ADMIN_ALLOWED_WALLETS + ADMIN_SECRET_KEY · POST /api/admin/auth",
    "• Maintenance mode — platform-wide freeze",
    "• Pause Crash engine independently · Pause Dice engine independently (503 on game APIs)",
    "• Kelly metrics & treasury health surfaces",
    "• Audit log of operator actions",
    "• Player / metrics views via GET /api/admin/metrics",
    "• Controls: POST /api/admin/controls → lib/admin-controls-state.ts",
    "• Complements circuit breaker for financial safety",
    "• Role: human-in-the-loop for incidents, migrations, and bankroll stress",
], size=15, color=WHITE, spacing=10)

# ═══════════════════════════════════════════════════════════
# SLIDE 21 — API & Data
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
section_header(s, "API Surface & Data Model", "Next.js App Router handlers + Supabase Postgres")
add_footer(s, 21)

add_rect(s, Inches(0.35), Inches(1.2), Inches(6.3), Inches(5.5), CARD, SLATE)
add_text(s, Inches(0.5), Inches(1.35), Inches(6), Inches(0.35),
         "Key Endpoints", size=15, bold=True, color=GOLD)
add_bullets(s, Inches(0.5), Inches(1.8), Inches(6), Inches(4.7), [
    "Auth: /nonce · /verify · /me · /logout",
    "User: /api/user/profile",
    "Dice: POST /api/games/dice/roll",
    "Crash: /bet · /cashout · /state · /stream · /ping",
    "PF: /provably-fair/verify · /rotate-seed",
    "VRF: /vrf/request · /vrf/verify",
    "Cashier: /deposit · /verify-deposit · /withdraw",
    "Treasury: /status · /stake",
    "Social: /trollbox · /bets/recent",
    "Security: /pow · /aml-check",
    "Admin: /auth · /controls · /metrics",
], size=12, color=WHITE, spacing=5)

add_rect(s, Inches(6.85), Inches(1.2), Inches(6.1), Inches(5.5), CARD, SLATE)
add_text(s, Inches(7.0), Inches(1.35), Inches(5.8), Inches(0.35),
         "Tables & RPCs", size=15, bold=True, color=VIOLET)
add_bullets(s, Inches(7.0), Inches(1.8), Inches(5.8), Inches(4.7), [
    "Tables: profiles · bets · transactions",
    "  trollbox_messages · bankroll_stakes",
    "  crash_rounds · aml_sanctions_quarantine",
    "",
    "RPCs (row locks):",
    "  execute_atomic_bet",
    "  execute_crash_wager",
    "  execute_crash_cashout",
    "",
    "RLS: public read on profiles/bets/trollbox/",
    "  stakes/crash_rounds; writes via service_role",
    "",
    "Fallback: in-memory mock if Supabase missing",
    "Project: cypherroll-casino",
], size=12, color=WHITE, spacing=4)

# ═══════════════════════════════════════════════════════════
# SLIDE 22 — Security invariants
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
section_header(s, "Security Invariants (Rules.md)", "Non-negotiable engineering constraints")
add_footer(s, 22)

inv = [
    ("No Math.random()", "All outcomes from HMAC-SHA256 / CSPRNG. Test suite enforces distribution."),
    ("Commit before play", "serverSeedHash published before bets. Reveal only after rotation / round end."),
    ("Never leak secrets", "No seeds, keys, or SESSION_SECRET in NEXT_PUBLIC_*. Client only sees hashes."),
    ("Atomic balances", "SQL FOR UPDATE / conditional .gte(balance, wager) prevents double-spend races."),
    ("Signed sessions", "cypher_session = base64url(payload).HMAC(SESSION_SECRET, payload)"),
    ("Wallet binding", "Game APIs require valid session in Real mode; spoofed addresses rejected."),
    ("Financial caps", "Kelly 1% · withdraw max $25k · solvency Σ balances ≤ reserves · circuit breaker."),
    ("Ops kill switches", "Maintenance + per-engine pause from admin — 503 on game routes."),
]
for i, (title, desc) in enumerate(inv):
    col, row = i % 2, i // 2
    x, y = Inches(0.4 + col * 6.45), Inches(1.2 + row * 1.35)
    add_rect(s, x, y, Inches(6.25), Inches(1.2), CARD, SLATE)
    add_text(s, x + Inches(0.2), y + Inches(0.15), Inches(5.8), Inches(0.3),
             title, size=14, bold=True, color=GOLD)
    add_text(s, x + Inches(0.2), y + Inches(0.55), Inches(5.8), Inches(0.5),
             desc, size=12, color=WHITE)

# ═══════════════════════════════════════════════════════════
# SLIDE 23 — Deployment
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
section_header(s, "Deployment Topology", "Docker · nginx · Tor · env surface")
add_footer(s, 23)

add_rect(s, Inches(0.4), Inches(1.25), Inches(6.2), Inches(5.5), CARD, SLATE)
add_text(s, Inches(0.6), Inches(1.4), Inches(5.8), Inches(0.35),
         "Runtime Path", size=16, bold=True, color=GOLD)
add_bullets(s, Inches(0.6), Inches(1.95), Inches(5.8), Inches(4.5), [
    "• Node 20 Alpine container (Next standalone)",
    "• BUILD_STANDALONE for Docker image",
    "• nginx reverse proxy on :8080",
    "• Tor v3 maps :80 → nginx",
    "• docker-compose wires app + env",
    "• Optional Redis for crash leader election",
    "• Supabase cloud or self-hosted PG",
    "• Public clearnet + .onion dual access",
], size=14, color=WHITE, spacing=8)

add_rect(s, Inches(6.85), Inches(1.25), Inches(6.0), Inches(5.5), CARD, SLATE)
add_text(s, Inches(7.05), Inches(1.4), Inches(5.6), Inches(0.35),
         "Critical Env Vars", size=16, bold=True, color=VIOLET)
add_bullets(s, Inches(7.05), Inches(1.95), Inches(5.6), Inches(4.5), [
    "• NEXT_PUBLIC_SUPABASE_URL / ANON_KEY",
    "• SUPABASE_SERVICE_ROLE_KEY",
    "• SESSION_SECRET · ADMIN_SECRET_KEY",
    "• ADMIN_ALLOWED_WALLETS",
    "• OPERATOR_SIGNER_PRIVATE_KEY",
    "• NEXT_PUBLIC_OPERATOR_SIGNER_ADDRESS",
    "• NEXT_PUBLIC_BASE_VAULT_ADDRESS",
    "• NEXT_PUBLIC_ARB_VAULT_ADDRESS",
    "• REDIS_URL / UPSTASH_*",
    "• POW_SECRET",
    "• ALLOW_SIMULATED_DEPOSITS (dev only)",
], size=12, color=WHITE, spacing=5)

# ═══════════════════════════════════════════════════════════
# SLIDE 24 — Demo vs Real
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
section_header(s, "Demo Mode vs Real Mode", "Same math — different money rails")
add_footer(s, 24)

add_rect(s, Inches(0.4), Inches(1.3), Inches(6.2), Inches(5.4), CARD, SLATE)
add_text(s, Inches(0.6), Inches(1.5), Inches(5.8), Inches(0.4),
         "DEMO (Guest Default)", size=18, bold=True, color=AMBER)
add_bullets(s, Inches(0.6), Inches(2.2), Inches(5.8), Inches(4.2), [
    "• No wallet required",
    "• Virtual $1,000 credits (localStorage)",
    "• Reset Credits available",
    "• Full PF math still runs",
    "• Demo roll reveals server seed immediately",
    "• No rakeback accrual",
    "• No on-chain settlement",
    "• Ideal for fairness skeptics to audit first",
], size=14, color=WHITE, spacing=8)

add_rect(s, Inches(6.85), Inches(1.3), Inches(6.0), Inches(5.4), CARD, SLATE)
add_text(s, Inches(7.05), Inches(1.5), Inches(5.6), Inches(0.4),
         "REAL (Authenticated)", size=18, bold=True, color=GREEN)
add_bullets(s, Inches(7.05), Inches(2.2), Inches(5.6), Inches(4.2), [
    "• Wallet + signed session required",
    "• USDC balance from vault deposits",
    "• Atomic SQL bet RPCs",
    "• Seeds stay secret until rotation",
    "• Rakeback accrues by VIP tier logic",
    "• Live ticker + social participation",
    "• Withdrawals via signed vouchers",
    "• Seamless switch from Demo after auth",
], size=14, color=WHITE, spacing=8)

# ═══════════════════════════════════════════════════════════
# SLIDE 25 — Status vs Roadmap
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
section_header(s, "Implementation Status vs Roadmap", "Honest inventory — prototype vs live mainnet")
add_footer(s, 25)

rows = [
    ("P0 Auth, PF, Dice, Crash, vaults, 3D", "Implemented", GREEN),
    ("P1 VIP, live bets, Kelly, Tor packaging", "Implemented", GREEN),
    ("P2 Community Bankroll LP", "UI + API/DB present", AMBER),
    ("P2 $CYPH buyback & burn", "Spec only", MUTED),
    ("Plinko / Martingale auto-bet", "PRD / incomplete", MUTED),
    ("Chainlink VRF live fulfillment", "Contract + simulated TS", AMBER),
    ("Vault addresses / Solana program", "Placeholders / devnet", AMBER),
    ("Treasury reserve figures", "Display defaults in UI", AMBER),
    ("Phase.docs (6 phases)", "Marked complete & verified", GREEN),
]
for i, (item, status, color) in enumerate(rows):
    y = Inches(1.2 + i * 0.6)
    add_rect(s, Inches(0.5), y, Inches(8.5), Inches(0.5), CARD, SLATE)
    add_rect(s, Inches(9.2), y, Inches(3.6), Inches(0.5), CARD, color)
    add_text(s, Inches(0.7), y + Inches(0.08), Inches(8.1), Inches(0.35),
             item, size=13, color=WHITE)
    add_text(s, Inches(9.2), y + Inches(0.08), Inches(3.6), Inches(0.35),
             status, size=12, bold=True, color=DARK if color != MUTED else WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 26 — User journeys
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
section_header(s, "End-to-End User Journeys", "Connect → Play → Settle → Audit")
add_footer(s, 26)

journeys = [
    ("Play Path", [
        "Connect wallet → Sign → Session",
        "Deposit USDC → AML verify → Credit",
        "Dice: set target → Roll → Settle",
        "Crash: bet in STARTING → Cash out",
        "Withdraw → EIP-712 / Solana claim",
    ]),
    ("Fairness Audit", [
        "Open Security / PF modal",
        "Note serverSeedHash pre-play",
        "Play → rotate or wait reveal",
        "Recompute HMAC locally",
        "Match outcome bit-for-bit",
    ]),
    ("LP Path", [
        "Open VAULT tab",
        "Check solvency % & health",
        "Stake USDC into bankroll",
        "Earn house-edge yield (APY UI)",
        "Monitor Kelly cap impact",
    ]),
    ("Ops Path", [
        "Operator wallet → /admin",
        "Review metrics & Kelly",
        "Pause engine if incident",
        "Toggle maintenance",
        "Clear quarantine / audit",
    ]),
]
for i, (title, items) in enumerate(journeys):
    x = Inches(0.35 + i * 3.25)
    add_rect(s, x, Inches(1.3), Inches(3.1), Inches(5.4), CARD, SLATE)
    add_text(s, x + Inches(0.15), Inches(1.45), Inches(2.8), Inches(0.4),
             title, size=15, bold=True, color=GOLD)
    add_bullets(s, x + Inches(0.15), Inches(2.1), Inches(2.8), Inches(4.3),
                [f"{j+1}. {t}" for j, t in enumerate(items)], size=12, color=WHITE, spacing=10)

# ═══════════════════════════════════════════════════════════
# SLIDE 27 — Key takeaways
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
section_header(s, "Key Takeaways", "What makes CypherRoll distinct")
add_footer(s, 27)

takes = [
    ("Hybrid settlement", "Off-chain HMAC games for speed; on-chain vaults for funds. Best of both worlds."),
    ("Math over trust", "Commit–reveal + public verifier turns fairness into a reproducible computation."),
    ("Risk engineered in", "Kelly 1%, $25k withdraw cap, solvency monitor, AML, PoW, admin kill switches."),
    ("Privacy as a product", "Tor v3 + zero-log posture + SIWE/SIWS zero-KYC is a first-class surface."),
    ("Rollbit-class UX", "3D canvases, live ticker, VIP, cashier, crash SSE — polished single-page casino lobby."),
    ("Clear honesty bar", "Rich prototype: PF + engines + schema + contracts. Placeholders on mainnet addrs / VRF / $CYPH."),
]
for i, (title, desc) in enumerate(takes):
    col, row = i % 3, i // 3
    x, y = Inches(0.4 + col * 4.25), Inches(1.3 + row * 2.7)
    add_rect(s, x, y, Inches(4.05), Inches(2.4), CARD, SLATE)
    add_text(s, x + Inches(0.2), y + Inches(0.25), Inches(3.6), Inches(0.4),
             title, size=15, bold=True, color=GOLD)
    add_text(s, x + Inches(0.2), y + Inches(0.8), Inches(3.6), Inches(1.3),
             desc, size=13, color=WHITE)

# ═══════════════════════════════════════════════════════════
# SLIDE 28 — Closing
# ═══════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
set_slide_bg(s)
add_sharp_rect(s, 0, Inches(2.6), prs.slide_width, Inches(0.08), GOLD)
add_text(s, Inches(0.8), Inches(1.8), Inches(11.5), Inches(0.5),
         "CYPHERROLL", size=48, bold=True, color=GOLD, align=PP_ALIGN.CENTER, font="Arial")
add_text(s, Inches(0.8), Inches(2.9), Inches(11.5), Inches(0.6),
         "Autonomous Provably Fair Gaming", size=22, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(3.7), Inches(11.5), Inches(0.8),
         "HMAC-SHA256  ·  Sub-50ms  ·  Multi-Chain Escrow  ·  Tor Privacy  ·  Kelly Solvency",
         size=14, color=AMBER, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(4.8), Inches(11.5), Inches(0.4),
         "Source: cypherroll-web3  ·  PRD · Architecture · Design · Rules · Phase.docs · lib/provably-fair.ts",
         size=12, color=MUTED, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.4),
         "“The best code is the code you never wrote.” — Ponytail Minimalist Engineering",
         size=13, color=VIOLET, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.3),
         "Thank you  ·  Questions welcome", size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

out = "/home/pratyush/Project/Website/cypherroll-web3/docs/CypherRoll_Technical_Deep_Dive.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
